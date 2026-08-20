"""Deterministic artifact comparators for file-oriented Office tasks."""

from __future__ import annotations

import csv
from collections import Counter
from dataclasses import dataclass
from difflib import SequenceMatcher
from pathlib import Path
from typing import Any

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation


def _normalize_text(text: str) -> str:
    return " ".join(text.split())


def _text_similarity(left: str, right: str) -> float:
    return SequenceMatcher(None, _normalize_text(left), _normalize_text(right)).ratio()


def compare_docx_files(expected_path: str, result_path: str, **options: Any) -> float:
    expected = Document(expected_path)
    result = Document(result_path)
    expected_text = "\n".join(paragraph.text for paragraph in expected.paragraphs if paragraph.text.strip())
    result_text = "\n".join(paragraph.text for paragraph in result.paragraphs if paragraph.text.strip())
    if options.get("ignore_case"):
        expected_text = expected_text.lower()
        result_text = result_text.lower()
    return 1.0 if _normalize_text(expected_text) == _normalize_text(result_text) else 0.0


def compare_docx_paragraph_styles(expected_path: str, result_path: str, **_: Any) -> float:
    expected = Document(expected_path)
    result = Document(result_path)
    expected_pairs = [(paragraph.text, paragraph.style.name) for paragraph in expected.paragraphs]
    result_pairs = [(paragraph.text, paragraph.style.name) for paragraph in result.paragraphs]
    return 1.0 if expected_pairs == result_pairs else 0.0


def compare_docx_tables(expected_path: str, result_path: str, **_: Any) -> float:
    expected = Document(expected_path)
    result = Document(result_path)
    if len(expected.tables) != len(result.tables):
        return 0.0
    for expected_table, result_table in zip(expected.tables, result.tables):
        if len(expected_table.rows) != len(result_table.rows):
            return 0.0
        if len(expected_table.columns) != len(result_table.columns):
            return 0.0
        for row_idx in range(len(expected_table.rows)):
            for col_idx in range(len(expected_table.columns)):
                left = expected_table.cell(row_idx, col_idx).text.strip()
                right = result_table.cell(row_idx, col_idx).text.strip()
                if left != right:
                    return 0.0
    return 1.0


def compare_docx_paragraph_slots(
    expected_path: str,
    result_path: str,
    *,
    expected_from_gold: list[int | dict[str, Any]],
    ignore_case: bool = False,
    **_: Any,
) -> float:
    expected = Document(expected_path)
    result = Document(result_path)

    matched = 0
    total = 0

    for item in expected_from_gold:
        spec = {"index": item} if isinstance(item, int) else dict(item)
        index = int(spec["index"]) - 1
        total += 1
        if index >= len(expected.paragraphs) or index >= len(result.paragraphs):
            continue

        expected_paragraph = expected.paragraphs[index]
        result_paragraph = result.paragraphs[index]
        expected_text = _normalize_text(expected_paragraph.text.lower() if ignore_case else expected_paragraph.text)
        result_text = _normalize_text(result_paragraph.text.lower() if ignore_case else result_paragraph.text)
        if expected_text != result_text:
            continue
        if spec.get("match_style") and expected_paragraph.style.name != result_paragraph.style.name:
            continue
        matched += 1

    if total == 0:
        return 0.0
    return round(matched / total, 4)


def compare_docx_expected_paragraphs(
    expected_path: str,
    result_path: str,
    *,
    expected_from_gold: list[int | dict[str, Any]],
    source_path: str | None = None,
    forbidden_from_source: list[int] | None = None,
    ignore_case: bool = False,
    **_: Any,
) -> float:
    expected = Document(expected_path)
    result = Document(result_path)
    source = Document(source_path) if source_path else None

    result_texts = [
        _normalize_text(paragraph.text.lower() if ignore_case else paragraph.text)
        for paragraph in result.paragraphs
        if paragraph.text.strip()
    ]
    result_pairs = [
        (
            _normalize_text(paragraph.text.lower() if ignore_case else paragraph.text),
            paragraph.style.name,
        )
        for paragraph in result.paragraphs
        if paragraph.text.strip()
    ]
    text_counter = Counter(result_texts)
    pair_counter = Counter(result_pairs)

    matched = 0
    total = 0

    for item in expected_from_gold:
        spec = {"index": item} if isinstance(item, int) else dict(item)
        paragraph = expected.paragraphs[int(spec["index"]) - 1]
        text = _normalize_text(paragraph.text.lower() if ignore_case else paragraph.text)
        total += 1
        if spec.get("match_style"):
            if pair_counter[(text, paragraph.style.name)] > 0:
                matched += 1
                pair_counter[(text, paragraph.style.name)] -= 1
        elif text_counter[text] > 0:
            matched += 1
            text_counter[text] -= 1

    if source is not None:
        for index in forbidden_from_source or []:
            paragraph = source.paragraphs[index - 1]
            if not paragraph.text.strip():
                continue
            total += 1
            text = _normalize_text(paragraph.text.lower() if ignore_case else paragraph.text)
            if text not in result_texts:
                matched += 1

    if total == 0:
        return 0.0
    return round(matched / total, 4)


def _paragraph_has_page_break(document: Document, paragraph_index1: int) -> bool:
    index = paragraph_index1 - 1
    if index < 0 or index >= len(document.paragraphs):
        return False

    paragraph_xml = document.paragraphs[index]._p.xml
    if "w:pageBreakBefore" in paragraph_xml:
        return True

    if index == 0:
        return False
    previous_xml = document.paragraphs[index - 1]._p.xml
    return 'w:type="page"' in previous_xml or "<w:lastRenderedPageBreak/>" in previous_xml


def compare_docx_expected_page_breaks(
    expected_path: str,
    result_path: str,
    *,
    expected_from_gold: list[int],
    **_: Any,
) -> float:
    expected = Document(expected_path)
    result = Document(result_path)

    total = len(expected_from_gold)
    if total == 0:
        return 0.0

    matched = 0
    for paragraph_index in expected_from_gold:
        if not _paragraph_has_page_break(expected, paragraph_index):
            continue
        if _paragraph_has_page_break(result, paragraph_index):
            matched += 1
    return round(matched / total, 4)


def _paragraph_has_drawing(document: Document, paragraph_index1: int) -> bool:
    index = paragraph_index1 - 1
    if index < 0 or index >= len(document.paragraphs):
        return False
    return "w:drawing" in document.paragraphs[index]._p.xml


def compare_docx_expected_drawings(
    expected_path: str,
    result_path: str,
    *,
    expected_from_gold: list[int],
    source_path: str | None = None,
    forbidden_from_source: list[int] | None = None,
    **_: Any,
) -> float:
    expected = Document(expected_path)
    result = Document(result_path)
    source = Document(source_path) if source_path else None

    matched = 0
    total = 0

    for paragraph_index in expected_from_gold:
        total += 1
        if _paragraph_has_drawing(expected, paragraph_index) and _paragraph_has_drawing(result, paragraph_index):
            matched += 1

    if source is not None:
        for paragraph_index in forbidden_from_source or []:
            total += 1
            if _paragraph_has_drawing(source, paragraph_index) and not _paragraph_has_drawing(result, paragraph_index):
                matched += 1

    if total == 0:
        return 0.0
    return round(matched / total, 4)


def _iter_slide_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(_normalize_text(shape.text))
    return texts


def compare_pptx_files(expected_path: str, result_path: str, **options: Any) -> float:
    expected = Presentation(expected_path)
    result = Presentation(result_path)
    expected_slides = list(expected.slides)
    result_slides = list(result.slides)
    if len(expected_slides) != len(result_slides):
        return 0.0

    expected_text = "\n".join("\n".join(_iter_slide_texts(slide)) for slide in expected_slides)
    result_text = "\n".join("\n".join(_iter_slide_texts(slide)) for slide in result_slides)
    if options.get("approximate") or options.get("approximately_tolerance") is not None:
        return round(_text_similarity(expected_text, result_text), 4)
    return 1.0 if expected_text == result_text else 0.0


def _slide_texts(presentation: Presentation, slide_numbers: list[int]) -> list[str]:
    slides = list(presentation.slides)
    texts: list[str] = []
    for slide_number in slide_numbers:
        slide = slides[slide_number - 1]
        texts.append("\n".join(_iter_slide_texts(slide)))
    return texts


def _relative_progress(source_score: float, result_score: float) -> float:
    if source_score >= 0.9999:
        return 1.0 if result_score >= 0.9999 else 0.0
    return max(0.0, min(1.0, (result_score - source_score) / (1.0 - source_score)))


def compare_pptx_changed_content(
    expected_path: str,
    result_path: str,
    *,
    source_path: str,
    changed_slides: list[int],
    unchanged_slides: list[int] | None = None,
    **_: Any,
) -> float:
    expected = Presentation(expected_path)
    result = Presentation(result_path)
    source = Presentation(source_path)

    expected_slides = list(expected.slides)
    result_slides = list(result.slides)
    source_slides = list(source.slides)
    if len(expected_slides) != len(result_slides) or len(source_slides) != len(result_slides):
        return 0.0

    changed_scores: list[float] = []

    for expected_text, result_text, source_text in zip(
        _slide_texts(expected, changed_slides),
        _slide_texts(result, changed_slides),
        _slide_texts(source, changed_slides),
    ):
        changed_scores.append(
            _relative_progress(
                _text_similarity(expected_text, source_text),
                _text_similarity(expected_text, result_text),
            )
        )

    unchanged_scores: list[float] = []
    for result_text, source_text in zip(
        _slide_texts(result, unchanged_slides or []),
        _slide_texts(source, unchanged_slides or []),
    ):
        unchanged_scores.append(1.0 if _normalize_text(result_text) == _normalize_text(source_text) else 0.0)

    if not changed_scores:
        return 0.0
    changed_score = sum(changed_scores) / len(changed_scores)
    unchanged_score = sum(unchanged_scores) / len(unchanged_scores) if unchanged_scores else 1.0
    return round(changed_score * unchanged_score, 4)


def compare_pdfs(expected_path: str, result_path: str, **_: Any) -> float:
    expected_reader = PdfReader(expected_path)
    result_reader = PdfReader(result_path)
    if len(expected_reader.pages) != len(result_reader.pages):
        return 0.0
    expected_text = "\n".join((page.extract_text() or "") for page in expected_reader.pages)
    result_text = "\n".join((page.extract_text() or "") for page in result_reader.pages)
    return round(_text_similarity(expected_text, result_text), 4)


def compare_csv(expected_path: str, result_path: str, **_: Any) -> float:
    def _read_rows(path: str) -> list[list[str]]:
        with open(path, newline="", encoding="utf-8") as handle:
            reader = csv.reader(handle)
            return [[cell.strip() for cell in row] for row in reader]

    return 1.0 if _read_rows(expected_path) == _read_rows(result_path) else 0.0


def compare_pdf_layout(expected_path: str, result_path: str, **_: Any) -> float:
    expected_reader = PdfReader(expected_path)
    result_reader = PdfReader(result_path)
    if len(expected_reader.pages) != len(result_reader.pages):
        return 0.0

    for expected_page, result_page in zip(expected_reader.pages, result_reader.pages):
        expected_box = tuple(round(float(value), 2) for value in expected_page.mediabox)
        result_box = tuple(round(float(value), 2) for value in result_page.mediabox)
        if expected_box != result_box:
            return 0.0
    return 1.0


def _trim_matrix(matrix: list[list[Any]]) -> list[list[Any]]:
    while matrix and all(value in (None, "") for value in matrix[-1]):
        matrix.pop()
    if not matrix:
        return []
    width = max((len(row) for row in matrix), default=0)
    while width > 0 and all((row[width - 1] if width - 1 < len(row) else None) in (None, "") for row in matrix):
        width -= 1
    return [row[:width] for row in matrix]


def _normalize_cell(value: Any, precision: int) -> Any:
    if value in (None, ""):
        return None
    if isinstance(value, float):
        return round(value, precision)
    if isinstance(value, str):
        return value.strip()
    return value


def _sheet_matrix(worksheet, precision: int) -> list[list[Any]]:
    rows: list[list[Any]] = []
    for row in worksheet.iter_rows(values_only=True):
        rows.append([_normalize_cell(value, precision) for value in row])
    return _trim_matrix(rows)


def _resolve_sheet_ref(sheet_ref: str | int, *, result_names: list[str], expected_names: list[str]) -> tuple[str, str]:
    if isinstance(sheet_ref, int):
        return ("result", result_names[sheet_ref])
    if sheet_ref.startswith("RI"):
        return ("result", result_names[int(sheet_ref[2:])])
    if sheet_ref.startswith("RN"):
        return ("result", sheet_ref[2:])
    if sheet_ref.startswith("EI"):
        return ("expected", expected_names[int(sheet_ref[2:])])
    if sheet_ref.startswith("EN"):
        return ("expected", sheet_ref[2:])
    raise ValueError(f"Unsupported sheet ref: {sheet_ref}")


def _pick_sheet(sheet_ref: str | int, *, result_workbook, expected_workbook):
    book_name, sheet_name = _resolve_sheet_ref(
        sheet_ref,
        result_names=result_workbook.sheetnames,
        expected_names=expected_workbook.sheetnames,
    )
    workbook = result_workbook if book_name == "result" else expected_workbook
    return workbook[sheet_name]


def _chart_title(chart) -> str | None:
    try:
        return chart.title.tx.rich.p[0].r[0].t
    except Exception:
        return None


def _chart_series_ref(series) -> str | None:
    num_ref = getattr(getattr(series, "val", None), "numRef", None)
    return getattr(num_ref, "f", None)


def _chart_series_title(series) -> str | None:
    tx = getattr(series, "tx", None)
    if tx is None:
        return None
    str_ref = getattr(tx, "strRef", None)
    if str_ref is not None:
        return getattr(str_ref, "f", None)
    return getattr(tx, "v", None)


def _load_chart_info(worksheet, chart_props: list[str]) -> list[dict[str, Any]]:
    charts: list[dict[str, Any]] = []
    for chart in getattr(worksheet, "_charts", []):
        info: dict[str, Any] = {}
        if "type" in chart_props:
            info["type"] = getattr(chart, "tagname", None)
        if "title" in chart_props:
            info["title"] = _chart_title(chart)
        if "series_refs" in chart_props:
            info["series_refs"] = [_chart_series_ref(series) for series in chart.series]
        if "series_titles" in chart_props:
            info["series_titles"] = [_chart_series_title(series) for series in chart.series]
        charts.append(info)
    return charts


def _load_pivot_info(worksheet, pivot_props: list[str]) -> list[dict[str, Any]]:
    pivots: list[dict[str, Any]] = []
    for pivot in getattr(worksheet, "_pivots", []):
        info: dict[str, Any] = {}
        if "row_fields" in pivot_props:
            info["row_fields"] = [field.x for field in pivot.rowFields]
        if "col_fields" in pivot_props:
            info["col_fields"] = [field.x for field in pivot.colFields]
        if "filter" in pivot_props:
            info["filter"] = [field.fld for field in pivot.pageFields]
        if "data_fields" in pivot_props:
            info["data_fields"] = [
                (field.fld, field.name, field.subtotal, field.showDataAs)
                for field in pivot.dataFields
            ]
        pivots.append(info)
    return pivots


@dataclass
class _PrintSettings:
    title_rows: str | None
    title_cols: str | None
    print_area: str | None
    fit_to_width: int | None
    fit_to_height: int | None
    orientation: str | None
    paper_size: int | None


def _print_settings(worksheet) -> _PrintSettings:
    return _PrintSettings(
        title_rows=worksheet.print_title_rows,
        title_cols=worksheet.print_title_cols,
        print_area=str(worksheet.print_area or ""),
        fit_to_width=worksheet.page_setup.fitToWidth,
        fit_to_height=worksheet.page_setup.fitToHeight,
        orientation=worksheet.page_setup.orientation,
        paper_size=worksheet.page_setup.paperSize,
    )


def compare_table(expected_path: str, result_path: str, **options: Any) -> float:
    result_workbook = load_workbook(result_path)
    expected_workbook = load_workbook(expected_path)

    for rule in options.get("rules", []):
        if rule["type"] == "sheet_data":
            precision = rule.get("precision", 4)
            left = _pick_sheet(rule["sheet_idx0"], result_workbook=result_workbook, expected_workbook=expected_workbook)
            right = _pick_sheet(rule["sheet_idx1"], result_workbook=result_workbook, expected_workbook=expected_workbook)
            if _sheet_matrix(left, precision) != _sheet_matrix(right, precision):
                return 0.0
        elif rule["type"] == "sheet_print":
            left = _pick_sheet(rule["sheet_idx0"], result_workbook=result_workbook, expected_workbook=expected_workbook)
            right = _pick_sheet(rule["sheet_idx1"], result_workbook=result_workbook, expected_workbook=expected_workbook)
            if _print_settings(left) != _print_settings(right):
                return 0.0
        elif rule["type"] == "chart":
            left = _pick_sheet(rule["sheet_idx0"], result_workbook=result_workbook, expected_workbook=expected_workbook)
            right = _pick_sheet(rule["sheet_idx1"], result_workbook=result_workbook, expected_workbook=expected_workbook)
            if _load_chart_info(left, rule.get("chart_props", [])) != _load_chart_info(right, rule.get("chart_props", [])):
                return 0.0
        elif rule["type"] == "pivot_table":
            left = _pick_sheet(rule["sheet_idx0"], result_workbook=result_workbook, expected_workbook=expected_workbook)
            right = _pick_sheet(rule["sheet_idx1"], result_workbook=result_workbook, expected_workbook=expected_workbook)
            if _load_pivot_info(left, rule.get("pivot_props", [])) != _load_pivot_info(right, rule.get("pivot_props", [])):
                return 0.0
        else:
            raise ValueError(f"Unsupported office table rule: {rule['type']}")
    return 1.0


def compare_xlsx_changed_cells(
    expected_path: str,
    result_path: str,
    *,
    source_path: str,
    sheets: list[str],
    precision: int = 4,
    **_: Any,
) -> float:
    expected_workbook = load_workbook(expected_path)
    result_workbook = load_workbook(result_path)
    source_workbook = load_workbook(source_path)

    total = 0
    matched = 0

    for sheet_name in sheets:
        expected_sheet = expected_workbook[sheet_name]
        result_sheet = result_workbook[sheet_name]
        source_sheet = source_workbook[sheet_name]
        max_row = max(expected_sheet.max_row, result_sheet.max_row, source_sheet.max_row)
        max_col = max(expected_sheet.max_column, result_sheet.max_column, source_sheet.max_column)

        for row_idx in range(1, max_row + 1):
            for col_idx in range(1, max_col + 1):
                expected_value = _normalize_cell(expected_sheet.cell(row_idx, col_idx).value, precision)
                source_value = _normalize_cell(source_sheet.cell(row_idx, col_idx).value, precision)
                if source_value == expected_value:
                    continue
                total += 1
                result_value = _normalize_cell(result_sheet.cell(row_idx, col_idx).value, precision)
                if result_value == expected_value:
                    matched += 1

    if total == 0:
        return 0.0
    return round(matched / total, 4)


def compare_xlsx_changed_chart_props(
    expected_path: str,
    result_path: str,
    *,
    source_path: str,
    sheet_name: str,
    chart_props: list[str],
    **_: Any,
) -> float:
    expected_sheet = load_workbook(expected_path)[sheet_name]
    result_sheet = load_workbook(result_path)[sheet_name]
    source_sheet = load_workbook(source_path)[sheet_name]

    expected_charts = _load_chart_info(expected_sheet, chart_props)
    result_charts = _load_chart_info(result_sheet, chart_props)
    source_charts = _load_chart_info(source_sheet, chart_props)
    if len(expected_charts) != len(result_charts) or len(expected_charts) != len(source_charts):
        return 0.0

    total = 0
    matched = 0
    for source_chart, expected_chart, result_chart in zip(source_charts, expected_charts, result_charts):
        for prop in chart_props:
            if source_chart.get(prop) == expected_chart.get(prop):
                continue
            total += 1
            if result_chart.get(prop) == expected_chart.get(prop):
                matched += 1

    if total == 0:
        return 0.0
    return round(matched / total, 4)


def compare_xlsx_expected_sheets(
    expected_path: str,
    result_path: str,
    *,
    sheets: list[str],
    precision: int = 4,
    **_: Any,
) -> float:
    expected_workbook = load_workbook(expected_path)
    result_workbook = load_workbook(result_path)

    total = len(sheets)
    if total == 0:
        return 0.0

    matched = 0
    for sheet_name in sheets:
        if sheet_name not in expected_workbook.sheetnames or sheet_name not in result_workbook.sheetnames:
            continue
        expected_sheet = expected_workbook[sheet_name]
        result_sheet = result_workbook[sheet_name]
        if _sheet_matrix(expected_sheet, precision) == _sheet_matrix(result_sheet, precision):
            matched += 1
    return round(matched / total, 4)


CHECK_FUNCTIONS = {
    "compare_csv": compare_csv,
    "compare_docx_files": compare_docx_files,
    "compare_docx_expected_drawings": compare_docx_expected_drawings,
    "compare_docx_expected_page_breaks": compare_docx_expected_page_breaks,
    "compare_docx_expected_paragraphs": compare_docx_expected_paragraphs,
    "compare_docx_paragraph_slots": compare_docx_paragraph_slots,
    "compare_docx_paragraph_styles": compare_docx_paragraph_styles,
    "compare_docx_tables": compare_docx_tables,
    "compare_pptx_files": compare_pptx_files,
    "compare_pptx_changed_content": compare_pptx_changed_content,
    "compare_pdfs": compare_pdfs,
    "compare_pdf_layout": compare_pdf_layout,
    "compare_table": compare_table,
    "compare_xlsx_changed_cells": compare_xlsx_changed_cells,
    "compare_xlsx_changed_chart_props": compare_xlsx_changed_chart_props,
    "compare_xlsx_expected_sheets": compare_xlsx_expected_sheets,
}


def run_check(func_name: str, expected_path: str, result_path: str, options: dict[str, Any] | None = None) -> float:
    try:
        func = CHECK_FUNCTIONS[func_name]
    except KeyError as exc:
        raise ValueError(f"Unsupported office artifact check: {func_name}") from exc
    return float(func(expected_path, result_path, **(options or {})))
